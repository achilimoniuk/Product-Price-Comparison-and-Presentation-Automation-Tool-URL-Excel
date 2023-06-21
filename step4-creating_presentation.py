import pandas as pd
import matplotlib.pyplot as plt
from matplotlib import pyplot as plt
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
import numpy as np
import dataframe_image as dfi
import time
from datetime import date
import seaborn as sns
import statistics
from tqdm import tqdm
from time import sleep
from termcolor import colored


# Funciton for creating statistics
def statfun(data, file, country):
    
    # Reading file with checked prices
    df = pd.read_csv(f'{data}.csv')

    # Choosing excel file
    excel = pd.read_excel(f'{file}.xlsx', header=0)
    excel = pd.DataFrame(excel)

    # Checking if contract id and program name is the same
    df['Contract ID from Excel'].replace('None', 0, inplace=True)
    df['Contract ID from URL'].replace('None', 0, inplace=True)
    df['Is Contract ID the same?'] = ''
    df['Is Price Program Name the same?'] = ''
    for index in df.index:
        if pd.to_numeric(df['Contract ID from URL'][index][-6:]) == df['Contract ID from Excel'][index]:
            df['Is Contract ID the same?'][index] = 'yes'
        else:
            df['Is Contract ID the same?'][index] = 'no'
        if df['Price Program from URL'][index] == df['Price Program from Excel'][index]:
            df['Is Price Program Name the same?'][index] = 'yes'
        else:
            df['Is Price Program Name the same?'][index] = 'no'

    # Saving different contracts and programs
    df_contractid = df.loc[(df['Is Contract ID the same?'] == 'no')]
    df_contractid.to_csv('contract different.csv')

    df_program = df.loc[(df['Is Price Program Name the same?'] == 'no')]
    df_program.to_csv('program different.csv')

    # Creating pie plot- percentage of different values
    def my_fmt(x):
        print(x)
        return '{:.2f}%\n({:.0f})'.format(x, total*x/100)
    counts = df['Is the same'].value_counts(sort=False)
    total = len(df['Is the same'])
    plt.pie(counts, labels = counts.index, autopct=my_fmt, colors = ['green','red'])
    plt.title('Are the prices the same?')
    plt.savefig("pieplot.png", bbox_inches='tight')
    plt.clf()

    # Barplot based on type of different values
    df['Difference type'].value_counts(sort=False).plot.bar(color=['green','red'])
    plt.ylabel("number of prices")
    plt.savefig("barplot.png", bbox_inches='tight')
    plt.clf()

    # Barplot based on contract ifd
    df['Is Contract ID the same?'].value_counts(sort=False).plot.bar(color=['green','red'])
    plt.ylabel("number of transactions")
    plt.title('Is the Contract ID the same?')
    plt.savefig("barplot contract.png", bbox_inches='tight')
    plt.clf()

    # Barplot based on program
    df['Is Price Program Name the same?'].value_counts(sort=False).plot.bar(color=['green','red'])
    plt.ylabel("number of transactions")
    plt.title('Is the Price Program Name the same?')
    plt.savefig("barplot program.png", bbox_inches='tight')
    plt.clf()

    # Table with numbers
    col = ['Name', 'Value']
    val1= ['Country',
    'Type of transaction',
    'Total # of Historical Sales Lines',
    'Total # of checked Sales Lines',
    '# of Lines with Successful Lookup',
    '# of Lines with Failed Lookup',
    '# of Lines with Failed Lookup- different values',
    '# of Lines with Failed Lookup- none value']
    val2 = [country, 
    df['Type of transaction'][0],
    excel.shape[0], 
    df.shape[0],  
    df[df['Difference type'] == 'no difference'].shape[0], 
    df.loc[df['Difference type'] !='no difference'].shape[0], 
    df[df['Difference type'] == 'different values'].shape[0], 
    df[df['Difference type'] == 'different values- none value'].shape[0]
    ]
    numbers = pd.DataFrame(zip(val1, val2), columns = col)
    dfi.export(numbers,"numbers.png")
    
    # Exporting table of different values to csv
    df3 = df.loc[(df['Difference type'] != 'no difference')]
    df3 = df3.reset_index()   
    df3 = df3.drop(['Is the same','Reference number',
    'Difference type','Is Contract ID the same?','Is Price Program Name the same?',
     'Type of transaction','index','Unnamed: 0'], axis=1)
    df3.to_csv('mismatched prices.csv') 
    df3.to_excel('mismatched prices.xlsx') 
    df3 =  df3.iloc[[0, 5, 12, 20, 52]]
    dfi.export(df3,"sample.png")


    ## DIFFERENT VALUES (WITHOUT NONE VALUES)
    # Creating histogram of incorrect values
    df2 = df.loc[(df['Difference type'] == 'different values')]
    df2['difference'] = pd.to_numeric(df2['difference'])
    df2 = df2.reset_index()   
    plot2 = plt.hist(df2['difference'])
    plt.title('Before removing outliers')
    plt.xlabel("difference")
    plt.ylabel("number of mismatched prices")
    plt.savefig("histogram_before.png", bbox_inches='tight')
    plt.clf()

    # Quantiles - different values
    quantiles =  df2['difference'].quantile(q=[0.05,0.25,0.5,0.75,0.95,1])
    quantiles.to_csv('q_before.csv')
    quantiles = pd.read_csv('q_before.csv', header=0, names=['Quantile', 'Value'])
    dfi.export(quantiles,"quantiles_before.png")

    # Frequency table - different values
    grouped_freq = df2['difference'].value_counts(bins = 20).sort_index()
    grouped_freq.to_csv('gf_before.csv')
    grouped_freq = pd.read_csv('gf_before.csv', header=0, names=['Values', 'Frequency'])
    grouped_freq['Frequency (%)'] = round(100*(grouped_freq['Frequency'] / grouped_freq['Frequency'].sum()), 2)
    grouped_freq['Cumulative Frequency'] = grouped_freq['Frequency'].cumsum()
    grouped_freq['Cumulative Frequency (%)'] = round(100*(grouped_freq['Frequency'].cumsum() / grouped_freq['Frequency'].sum()), 2)
    dfi.export(grouped_freq,"frequency_before.png")

    # Box Plot- detecting outliers 
    sns.boxplot(df2['difference'])
    plt.savefig("boxplot.png", bbox_inches='tight')
    plt.clf()

    # Detecting outliers
    df3 = pd.DataFrame(df2)
    Q1 = np.percentile(df3['difference'], 25,
                    interpolation = 'midpoint')
    
    Q3 = np.percentile(df3['difference'], 75,
                    interpolation = 'midpoint')
    IQR = Q3 - Q1 
    # Upper bound
    upper = np.where(df3['difference'] >= (Q3+1.5*IQR))
    # Lower bound
    lower = np.where(df3['difference'] <= (Q1-1.5*IQR))
    # Merging lower and upper bound
    out = np.concatenate((upper[0], lower[0]))

    # Exporting the table of outliers to png
    outliers = df3.iloc[out]
    outliers2 = outliers.drop(['index', 'Unnamed: 0', 'Is the same', 
    'Difference type', 'Type of transaction', 'Is Contract ID the same?', 'Is Price Program Name the same?'], axis=1) 
    dfi.export(outliers2,"outliers.png")

    #Removing the Outliers 
    df3.drop(out, inplace = True)

    # Statistics
    stats = []
    stats.append(['mean', round(statistics.mean(df2['difference']),3), round(statistics.mean(df3['difference']),3)])
    stats.append(['median', round(statistics.median(df2['difference']),3), round(statistics.median(df3['difference']),3)])
    stats.append(['standard deviation', round(statistics.stdev(df2['difference']),3), round(statistics.stdev(df3['difference']),3)])
    stats.append(['minimum difference', min(df2['difference']), min(df3['difference'])])
    stats.append(['maximum difference', round(max(df2['difference']),3), round(max(df3['difference']),3)])
    stats = pd.DataFrame(stats, columns= ['Statisic', 'Before removing outliers', 'After removing outliers'])
    dfi.export(stats,"statistics.png")

    # Creating histogram of incorrect values- after removing the outliers
    plot3 = plt.hist(df3['difference'])
    plt.title('After removing outliers')
    plt.xlabel("difference")
    plt.ylabel("number of mismatched prices")
    plt.savefig("histogram_after.png", bbox_inches='tight')
    plt.clf()

    # Quantiles - different values- after removing the outliers
    quantiles2 =  df3['difference'].quantile(q=[0.05,0.25,0.5,0.75,0.95,1])
    quantiles2.to_csv('q_after.csv')
    quantiles2 = pd.read_csv('q_after.csv', header=0, names=['Quantile', 'Value'])
    dfi.export(quantiles2,"quantiles_after.png")

    # Frequency table - different values
    grouped_freq2 = df3['difference'].value_counts(bins = 20).sort_index()
    grouped_freq2.to_csv('gf_after.csv')
    grouped_freq2 = pd.read_csv('gf_after.csv', header=0, names=['Values', 'Frequency'])
    grouped_freq2['Frequency (%)'] = round(100*(grouped_freq2['Frequency'] / grouped_freq2['Frequency'].sum()), 2)
    grouped_freq2['Cumulative Frequency'] = grouped_freq2['Frequency'].cumsum()
    grouped_freq2['Cumulative Frequency (%)'] = round(100*(grouped_freq2['Frequency'].cumsum() / grouped_freq2['Frequency'].sum()), 2)
    dfi.export(grouped_freq2,"frequency_after.png")

    # DIFFERENT VALUES- NONE VALUE
    df4 = df.loc[(df['Difference type'] == 'different values- none value')]
    df4 = df4.reset_index() 
    df4 = df4.drop(['Is the same', 'Difference type', 'index', 'Unnamed: 0'], axis=1)  
    dfi.export(df4,"tablenone.png")


    # PPTX EXPORT
    # Opening presentation and creating first slide
    ppt = Presentation('presentation.pptx')
    first_slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    title = "Statistics  - " + str(date.today())
    first_slide.shapes.title.text = title

    # Creating titles for the slides
    title2 = "Details of historical and checked lines"
    title3 = "Number of matched and mismatched prices"
    title4 = "Number of matched and mismatched prices grouped by type of match and mismatch"
    title5 = "Boxplot of mismatched prices- detecting outliers"
    title6 = "Detected outliers"
    title7 = "Histogram of mismatched prices before and after removing outliers"
    title8 = "Quantiles of mismatched prices before and after removing outliers"
    title9 = "Frequency and cumulative frequency of mismatched prices before and after removing outliers"
    title10 = "Statistics before and after removing outliers"
    title11 = "Details of prices with none values in Abbott Service"
    title12 = "Number of matched and mismatched Contract ID and Price Program Name"
   

    # Setting the objects
    num = 'numbers.png'
    pie = 'pieplot.png'
    bar = 'barplot.png'
    hist_b = 'histogram_before.png'
    quant_b = 'quantiles_before.png'
    freq_b = 'frequency_before.png'
    box = 'boxplot.png'
    outlr = 'outliers.png'
    stats = 'statistics.png'
    hist_a = 'histogram_after.png'
    quant_a = 'quantiles_after.png'
    freq_a = 'frequency_after.png'
    none = 'tablenone.png'
    contr = 'barplot contract.png'
    program = 'barplot program.png'

    # Creating slides
    slide2 = ppt.slides.add_slide(ppt.slide_layouts[1])
    slide3 = ppt.slides.add_slide(ppt.slide_layouts[1])
    slide4 = ppt.slides.add_slide(ppt.slide_layouts[1])
    slide5 = ppt.slides.add_slide(ppt.slide_layouts[1])
    slide6 = ppt.slides.add_slide(ppt.slide_layouts[1])
    slide7 = ppt.slides.add_slide(ppt.slide_layouts[1])
    slide8 = ppt.slides.add_slide(ppt.slide_layouts[1])
    slide9 = ppt.slides.add_slide(ppt.slide_layouts[1])
    slide10 = ppt.slides.add_slide(ppt.slide_layouts[1])
    slide11 = ppt.slides.add_slide(ppt.slide_layouts[1])
    slide12 = ppt.slides.add_slide(ppt.slide_layouts[1])
    # Titles
    slide2.shapes.title.text = title2
    slide3.shapes.title.text = title3
    slide4.shapes.title.text = title4
    slide5.shapes.title.text = title5
    slide6.shapes.title.text = title6
    slide7.shapes.title.text = title7
    slide8.shapes.title.text = title8
    slide9.shapes.title.text = title9
    slide10.shapes.title.text = title10
    slide11.shapes.title.text = title11
    slide12.shapes.title.text = title12
   

    # Pasting charts and tables into the slides
    num = slide2.shapes.add_picture(num, left= Inches(3),top = Inches(2),height = Inches(4))
    pie = slide3.shapes.add_picture(pie, left= Inches(3),top = Inches(3),height = Inches(5))
    bar = slide4.shapes.add_picture(bar, left= Inches(2),top = Inches(3),height = Inches(5))
    box = slide5.shapes.add_picture(box, left= Inches(3),top = Inches(2),height = Inches(5))
    outlr = slide6.shapes.add_picture(outlr, left= Inches(3),top = Inches(2),height = Inches(5))
    hist = slide7.shapes.add_picture(hist_b, left= Inches(1),top = Inches(3),height = Inches(4))
    hist = slide7.shapes.add_picture(hist_a, left= Inches(5),top = Inches(3),height = Inches(4))
    quant = slide8.shapes.add_picture(quant_b, left= Inches(1),top = Inches(2),height = Inches(4))
    quant = slide8.shapes.add_picture(quant_a, left= Inches(5),top = Inches(2),height = Inches(4))
    freq = slide9.shapes.add_picture(freq_b, left= Inches(1),top = Inches(2),height = Inches(4))
    freq = slide9.shapes.add_picture(freq_a, left= Inches(6),top = Inches(2),height = Inches(4))
    stats = slide10.shapes.add_picture(stats, left= Inches(2),top = Inches(2),height = Inches(3))
    none = slide11.shapes.add_picture(none, left= Inches(2),top = Inches(2),height = Inches(4))
    contr = slide12.shapes.add_picture(contr, left= Inches(1),top = Inches(3),height = Inches(4))
    program = slide12.shapes.add_picture(program, left= Inches(5),top = Inches(3),height = Inches(4))
  
    # Saving the powerpoint presentation
    ppt.save('presentation.pptx')

# Going through the files
files=['file1', 'file2']
for i in tqdm(files, desc = colored("Update on creating statistics:", "red"), colour="green"):
    statfun('prices done', files[i], 'conutry')
print(colored("Statistics created", "green"))
print(colored("Powerpoint presentation exported", "green"))

